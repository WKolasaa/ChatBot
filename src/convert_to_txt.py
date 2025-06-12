import os
from docx import Document
import pdfplumber
from pptx import Presentation

RAW_DIR = "data/raw"
TXT_DIR = "data/raw_txt"
os.makedirs(TXT_DIR, exist_ok=True)

def convert_docx_to_txt(docx_path, txt_path):
    doc = Document(docx_path)
    text = "\n".join([para.text for para in doc.paragraphs])
    with open(txt_path, "w", encoding="utf-8") as f:
        f.write(text)

def convert_pdf_to_txt(pdf_path, txt_path):
    with pdfplumber.open(pdf_path) as pdf:
        pages = [page.extract_text() or "" for page in pdf.pages]
    text = "\n".join(pages)
    with open(txt_path, "w", encoding="utf-8") as f:
        f.write(text)

def convert_pptx_to_txt(pptx_path, txt_path):
    prs = Presentation(pptx_path)
    texts = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                texts.append(shape.text)
    text = "\n".join(texts)
    with open(txt_path, "w", encoding="utf-8") as f:
        f.write(text)

def convert_all_files():
    for fname in os.listdir(RAW_DIR):
        ext = fname.lower().split(".")[-1]
        in_path = os.path.join(RAW_DIR, fname)
        out_fname = fname + ".txt"  # Ensures uniqueness
        out_path = os.path.join(TXT_DIR, out_fname)
        try:
            if ext == "docx":
                convert_docx_to_txt(in_path, out_path)
                print(f"Converted {fname} to txt.")
            elif ext == "pdf":
                convert_pdf_to_txt(in_path, out_path)
                print(f"Converted {fname} to txt.")
            elif ext == "pptx":
                convert_pptx_to_txt(in_path, out_path)
                print(f"Converted {fname} to txt.")
            elif ext == "txt":
                # Already text, just copy
                with open(in_path, "r", encoding="utf-8") as src, open(out_path, "w", encoding="utf-8") as dst:
                    dst.write(src.read())
                print(f"Copied {fname}.")
            else:
                print(f"Skipped unsupported file: {fname}")
        except Exception as e:
            print(f"Failed to process {fname}: {e}")

if __name__ == "__main__":
    convert_all_files()
